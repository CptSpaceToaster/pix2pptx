#!/usr/bin/env python3.4
# apt-get install libxml2-dev libxslt1-dev
# pip3 install python-pptx
from pptx import Presentation
from pptx.util import Px
# pip3 install pillow
from PIL import Image
# pip3 install naturalsort
from natsort import natsorted, ns
import argparse
import os

parser = argparse.ArgumentParser(description='pix2pptx - Convert a folder of images into a power point presentation')
mandatory = parser.add_argument_group('mandatory arguments')
# TODO: List what type of image files are permitted
mandatory.add_argument('dir', help='Directory location containing image files')
parser.add_argument('-f', '--filename', default='pix.pptx',
                    help='Output filename (default="out.pptx")')
args = parser.parse_args()

# Since the default presentation is 10″ x 7.5″ the size of each page is 720 x 540.
slide_width = 720
slide_height = 540

if __name__ == '__main__':
    print(args.dir)
    P = Presentation()

    if os.path.isdir(args.dir):
        for root, dirs, files in os.walk(args.dir):
            # using natsort>=3.5.0
            for img in natsorted(files, alg=ns.IGNORECASE):
                img = root + os.sep + img

                # get image size in pixels
                cat = Image.open(img)
                img_width = cat.size[0]
                img_height = cat.size[1]

                scale = min(slide_width/img_width, slide_height/img_height)

                # calculate offsets
                left_offset = (slide_width-(img_width*scale))/2
                top_offset = (slide_height-(img_height*scale))/2

                # add a new slide
                slide = P.slides.add_slide(P.slide_layouts[6])

                # add image to new slide
                pic = slide.shapes.add_picture(img, Px(left_offset), Px(top_offset), width=Px(img_width*scale), height=Px(img_height*scale))

    # we're done here
    P.save(args.dir + os.sep + args.filename)
